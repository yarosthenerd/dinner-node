from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def C(t): return RGBColor(int(t[0:2],16), int(t[2:4],16), int(t[4:6],16))
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

def slide(items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = C('0B0E14')
    y = 1.0
    for it in items:
        text, color, size, bold = it[:4]
        url = it[4] if len(it) > 4 else None
        tb = s.shapes.add_textbox(Inches(0.9), Inches(y), Inches(11.5), Inches(2.2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = 'Consolas'
        r.font.color.rgb = C(color)
        if url:
            r.hyperlink.address = url; r.font.underline = True
        y += 0.4 + (size / 72.0) * 1.7
    return s

slide([('DinnerNode \u25cf','9FEF00',54,True),
       ('your idle PC pays for dinner.','9FEF00',30,False),
       ('rent the world\u2019s idle hardware \u00b7 pay per second \u00b7 on Monad','6B7A89',18,False)])

slide([('the problem','9FEF00',40,True),
       ('privacy \u2014 every prompt is read, logged and monetized on someone else\u2019s server.','D7E0EA',22,False),
       ('stalled hardware \u2014 billions of GPUs idle ~90% of the time.','D7E0EA',22,False),
       ('AI subscriptions \u2014 $20\u2013200/month while your own machine earns nothing.','D7E0EA',22,False)])

slide([('the solution','9FEF00',40,True),
       ('rent a real machine directly \u00b7 pay per token, per second \u00b7 no middleman.','D7E0EA',22,False),
       ('your prompt stays between you and the machine you rented.','D7E0EA',22,False),
       ('why Monad: one answer = ~30 micropayments \u00b7 Ethereum \u2248 $115 \u00b7 Monad \u2248 $0.0003','9FEF00',22,False)])

slide([('try it live','9FEF00',44,True),
       ('web-opal-sigma-55.vercel.app','9FEF00',30,True,'https://web-opal-sigma-55.vercel.app/'),
       ('one tap: place order \u00b7 watch the check tick \u00b7 every token is a tip.','6B7A89',20,False)])

prs.save('DinnerNode_pitch.pptx')
print('saved DinnerNode_pitch.pptx')
